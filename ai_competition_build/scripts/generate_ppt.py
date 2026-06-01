from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
AI = next(p for p in ROOT.iterdir() if p.is_dir() and p.name.startswith("ai") and "竞赛资料" in p.name)
OUT = AI / "AI竞赛综合汇报_大图连续浏览SDK与展示平台.pptx"
PREVIEW_DIR = AI / "AI竞赛综合汇报_PPT预览图"
PREVIEW_DIR.mkdir(exist_ok=True)
LOGO = AI / "汉宁logo.png"


def find_dir(base: Path, *parts: str) -> Path:
    current = base
    for part in parts:
        current = next(p for p in current.iterdir() if p.is_dir() and part in p.name)
    return current


SDK_IMG = sorted(find_dir(AI, "大图连续浏览SDK参赛交付包", "软件界面图像").glob("*.png"))
PLATFORM_IMG = sorted(find_dir(AI, "展示平台参赛交付包", "软件界面图像").glob("*.png"))

C = {
    "ink": RGBColor(13, 27, 42),
    "muted": RGBColor(100, 116, 139),
    "line": RGBColor(226, 232, 240),
    "blue": RGBColor(30, 58, 95),
    "cyan": RGBColor(56, 189, 248),
    "green": RGBColor(34, 197, 94),
    "yellow": RGBColor(255, 163, 0),
    "red": RGBColor(244, 114, 24),
    "navy": RGBColor(16, 36, 56),
    "dark": RGBColor(13, 27, 42),
    "pale_blue": RGBColor(241, 245, 249),
    "pale_green": RGBColor(240, 253, 244),
    "white": RGBColor(255, 255, 255),
    "soft": RGBColor(248, 250, 252),
}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_text(slide, text, x, y, w, h, size=12, color=None, bold=False, align=None, valign=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if valign:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "微软雅黑"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C["ink"]
    return box


def add_multiline(slide, items, x, y, w, h, size=11, color=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(5)
        r = p.add_run()
        r.text = "• " + item
        r.font.name = "微软雅黑"
        r.font.size = Pt(size)
        r.font.color.rgb = color or C["ink"]
    return box


def add_panel(slide, x, y, w, h, fill=None, line=None):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill or C["white"]
    shp.line.color.rgb = line or C["line"]
    shp.line.width = Pt(0.8)
    return shp


def add_brand(slide, page_no):
    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(8.0), Inches(0.06))
    left.fill.solid()
    left.fill.fore_color.rgb = C["yellow"]
    left.line.color.rgb = C["yellow"]
    right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.0), Inches(0), Inches(5.33), Inches(0.06))
    right.fill.solid()
    right.fill.fore_color.rgb = C["cyan"]
    right.line.color.rgb = C["cyan"]
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(0.18), Inches(0.12), width=Inches(0.42))
    add_text(slide, f"{page_no:02d}", 12.15, 6.92, 0.55, 0.18, size=8, color=C["muted"], align=PP_ALIGN.RIGHT)


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.75, 0.2, 6.9, 0.5, size=23, color=C["blue"], bold=True, valign=False)
    if subtitle:
        add_text(slide, subtitle, 0.75, 0.7, 9.6, 0.35, size=10.5, color=C["muted"], valign=False)


def add_label(slide, label, x, y, color=None):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(1.2), Inches(0.28))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color or C["blue"]
    shp.line.color.rgb = color or C["blue"]
    add_text(slide, label, x + 0.08, y + 0.02, 1.04, 0.22, size=7.6, color=C["white"], bold=True, align=PP_ALIGN.CENTER)


def add_metric(slide, value, label, note, x, y, w, color):
    add_text(slide, value, x, y, w, 0.44, size=25, color=color, bold=True, valign=False)
    add_text(slide, label, x, y + 0.48, w, 0.22, size=8.5, color=C["ink"], bold=True, valign=False)
    add_text(slide, note, x, y + 0.72, w, 0.34, size=7.2, color=C["muted"], valign=False)


def add_placeholder(slide, x, y, w, h, label):
    add_panel(slide, x, y, w, h, C["soft"], RGBColor(184, 198, 211))
    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.05))
    strip.fill.solid()
    strip.fill.fore_color.rgb = C["yellow"]
    strip.line.color.rgb = C["yellow"]
    add_text(slide, label, x + 0.2, y + h / 2 - 0.15, w - 0.4, 0.3, size=12, color=C["muted"], align=PP_ALIGN.CENTER)


def add_image_contain(slide, img, x, y, w, h, caption=None):
    if img is None or not img.exists():
        add_placeholder(slide, x, y, w, h, caption or "待补充截图")
        return
    with Image.open(img) as im:
        iw, ih = im.size
    ratio = iw / ih
    ww = w
    hh = w / ratio
    if hh > h:
        hh = h
        ww = h * ratio
    slide.shapes.add_picture(str(img), Inches(x + (w - ww) / 2), Inches(y + (h - hh) / 2), width=Inches(ww), height=Inches(hh))
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.background()
    rect.line.color.rgb = C["line"]
    rect.line.transparency = 20
    if caption:
        add_text(slide, caption, x, y + h + 0.05, w, 0.24, size=7.5, color=C["muted"], align=PP_ALIGN.CENTER, valign=False)


def add_image_cover(slide, img, x, y, w, h):
    if img is None or not img.exists():
        add_placeholder(slide, x, y, w, h, "待补充截图")
        return
    with Image.open(img) as im:
        iw, ih = im.size
    ratio = iw / ih
    target = w / h
    if ratio > target:
        hh = h
        ww = h * ratio
    else:
        ww = w
        hh = w / ratio
    slide.shapes.add_picture(str(img), Inches(x + (w - ww) / 2), Inches(y + (h - hh) / 2), width=Inches(ww), height=Inches(hh))


def add_table(slide, rows, x, y, widths, row_h=0.42):
    for r_idx, row in enumerate(rows):
        cx = x
        for c_idx, cell in enumerate(row):
            fill = C["navy"] if r_idx == 0 else (C["white"] if r_idx % 2 else RGBColor(248, 250, 252))
            rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(cx), Inches(y + r_idx * row_h), Inches(widths[c_idx]), Inches(row_h))
            rect.fill.solid()
            rect.fill.fore_color.rgb = fill
            rect.line.color.rgb = RGBColor(223, 231, 239)
            rect.line.width = Pt(0.4)
            add_text(
                slide,
                str(cell),
                cx + 0.06,
                y + r_idx * row_h + 0.04,
                widths[c_idx] - 0.12,
                row_h - 0.08,
                size=7.5 if r_idx else 7.8,
                color=C["white"] if r_idx == 0 else C["ink"],
                bold=r_idx == 0,
                valign=True,
            )
            cx += widths[c_idx]


def add_bars(slide, rows, x, y, w, max_val, color_a, color_b):
    for i, (label, before, after) in enumerate(rows):
        top = y + i * 0.56
        add_text(slide, label, x, top, 2.2, 0.22, size=8, color=C["ink"], valign=False)
        bw = (w - 3.1) * before / max_val
        aw = (w - 3.1) * after / max_val
        b = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 2.35), Inches(top + 0.08), Inches(bw), Inches(0.12))
        b.fill.solid(); b.fill.fore_color.rgb = color_a; b.line.color.rgb = color_a
        a = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 2.35), Inches(top + 0.27), Inches(aw), Inches(0.12))
        a.fill.solid(); a.fill.fore_color.rgb = color_b; a.line.color.rgb = color_b
        add_text(slide, f"{before}h / {after}h", x + w - 0.72, top + 0.1, 0.7, 0.18, size=7.2, color=C["muted"], align=PP_ALIGN.RIGHT, valign=False)
    add_text(slide, "传统估算", x + 2.35, y + len(rows) * 0.56 + 0.02, 0.75, 0.16, size=7, color=color_a, valign=False)
    add_text(slide, "AI协作", x + 3.16, y + len(rows) * 0.56 + 0.02, 0.75, 0.16, size=7, color=color_b, valign=False)


def new_slide(bg="FFFFFF"):
    slide = prs.slides.add_slide(BLANK)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string(bg)
    return slide


def build():
    page = 1
    slide = new_slide("0D1B2A")
    dark = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(8.5), Inches(7.5))
    dark.fill.solid(); dark.fill.fore_color.rgb = C["dark"]; dark.line.color.rgb = C["dark"]
    add_image_cover(slide, SDK_IMG[3], 8.0, 0.0, 5.33, 7.5)
    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.0), Inches(0), Inches(5.33), Inches(7.5))
    overlay.fill.solid(); overlay.fill.fore_color.rgb = C["dark"]; overlay.fill.transparency = 36; overlay.line.color.rgb = C["dark"]
    add_brand(slide, page); page += 1
    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(0.82), Inches(3.2), Inches(0.38))
    tag.fill.solid(); tag.fill.fore_color.rgb = C["blue"]; tag.line.color.rgb = C["blue"]
    add_text(slide, "AI 应用竞赛 · 2026", 0.75, 0.86, 3.0, 0.28, size=11, color=C["white"], bold=True, valign=False)
    add_text(slide, "AI 赋能研发交付", 0.7, 1.82, 7.6, 0.62, size=28, color=C["white"], bold=True, valign=False)
    add_text(slide, "大图连续浏览 SDK 与地铁隧道展示平台", 0.7, 2.58, 7.8, 0.82, size=25, color=C["cyan"], bold=True, valign=False)
    add_text(slide, "按项目单元汇报 AI 工具、软件成果、量化指标与实践体会", 0.75, 3.55, 6.8, 0.36, size=13, color=RGBColor(226, 232, 240), valign=False)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.75), Inches(5.58), Inches(7.25), Inches(0.08))
    accent.fill.solid(); accent.fill.fore_color.rgb = C["yellow"]; accent.line.color.rgb = C["yellow"]
    add_text(slide, "汇报资料：README、竞赛说明、既有展示平台交付包、软件运行截图", 0.75, 5.92, 6.8, 0.32, size=10.5, color=RGBColor(226, 232, 240), valign=False)
    add_text(slide, "公司 AI 应用竞赛参赛材料", 0.75, 6.28, 6.0, 0.3, size=10, color=RGBColor(148, 163, 184), valign=False)

    slide = new_slide()
    add_brand(slide, page); page += 1
    add_title(slide, "目录 · Contents", "按项目单元组织，先讲 AI 过程，再讲软件成果、量化指标与体会。")
    cards = [
        ("01", "整体汇报主线", "两个项目串成从大图处理到成果展示的研发交付链路"),
        ("02", "SDK · AI 协作过程", "AI 工具、沟通过程、需求拆解和工程落地方式"),
        ("03", "SDK · 成果与指标", "LOD 浏览、高清瓦片、标注能力和代码/工时测算"),
        ("04", "展示平台 · AI 协作过程", "从编译失败、数据映射到可演示系统的修复闭环"),
        ("05", "展示平台 · 成果与指标", "上传工具、统一库、WebAPI、页面截图和量化看板"),
        ("06", "综合结论与推广", "风险治理、过程证据和可复制 AI 研发实践模板"),
    ]
    positions = [(0.5, 1.3), (4.8, 1.3), (9.1, 1.3), (0.5, 3.9), (4.8, 3.9), (9.1, 3.9)]
    for i, ((no, title, body), (x, y)) in enumerate(zip(cards, positions)):
        add_panel(slide, x, y, 4.0, 2.3, C["white"], C["line"])
        strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(4.0), Inches(0.05))
        strip.fill.solid(); strip.fill.fore_color.rgb = C["yellow"] if i % 2 == 0 else C["cyan"]; strip.line.color.rgb = strip.fill.fore_color.rgb
        add_text(slide, no, x + 0.15, y + 0.1, 0.7, 0.55, size=24, color=C["yellow"] if i % 2 == 0 else C["cyan"], bold=True, valign=False)
        add_text(slide, title, x + 0.15, y + 0.65, 3.65, 0.35, size=13, color=C["blue"], bold=True, valign=False)
        add_text(slide, body, x + 0.15, y + 1.1, 3.62, 0.55, size=9.6, color=C["muted"], valign=False)

    slide = new_slide()
    add_brand(slide, page); page += 1
    add_title(slide, "项目一：大图连续浏览 SDK - AI 工具与沟通过程", "AI 参与需求拆解、架构选型、代码生成、性能策略和问题修复，人工负责业务判断和验收。")
    add_label(slide, "AI工具", 0.75, 1.78, C["blue"])
    add_multiline(slide, [
        "ChatGPT / Codex：用于架构讨论、Qt/C++ 代码实现、README 与竞赛材料整理。",
        "OpenCode 千问、MiniMax 等：用于局部 Bug 思路、替代方案和表达优化。",
        "人工把关：以真实编译、运行截图、业务规则和工程约束作为最终验收标准。",
    ], 0.75, 2.18, 5.0, 1.05, size=11)
    add_label(slide, "沟通过程", 0.75, 3.65, C["green"])
    add_table(slide, [
        ["阶段", "AI协作内容", "形成结果"],
        ["场景澄清", "说明 GB 级隧道影像浏览痛点", "明确切片、LOD、异步加载路线"],
        ["模块实现", "拆分 SDK、DemoApp、ImageSlicer", "形成可运行工程结构"],
        ["迭代修复", "围绕崩溃、编码、64位构建等问题排查", "README 与使用说明沉淀"],
    ], 0.75, 4.05, [1.05, 2.25, 2.3], 0.43)
    add_placeholder(slide, 6.85, 1.82, 5.55, 4.75, "待补充：AI 对话过程截图 / Codex 修改记录 / 关键提示词截图")

    slide = new_slide("F8FBFD")
    add_brand(slide, page); page += 1
    add_title(slide, "项目一：软件介绍与技术亮点", "面向 GB 级隧道全景图，提供切片预处理、连续浏览、缩放、标注和数据保存能力。")
    add_image_contain(slide, SDK_IMG[2], 0.75, 1.82, 5.45, 3.08, "缩略图 / LOD 浏览界面")
    add_image_contain(slide, SDK_IMG[3], 6.65, 1.82, 5.45, 3.08, "高清瓦片加载运行效果")
    add_multiline(slide, [
        "Google Maps 式瓦片加载：只加载视野内切片，降低大图内存压力。",
        "LOD 机制：缩小时用缩略图快速定位，放大后异步加载高清切片。",
        "后台线程解码：图片读取与解码不阻塞主界面，保障浏览交互流畅。",
        "专业标注：支持折线、多边形、病害类型录入、右键管理和 JSON 保存。",
    ], 0.9, 5.38, 11.2, 0.92, size=10.2)

    slide = new_slide()
    add_brand(slide, page); page += 1
    add_title(slide, "项目一：运行效果截图", "从原始影像切片到 SDK 浏览，展示端到端可运行结果。")
    add_image_contain(slide, SDK_IMG[0], 0.75, 1.72, 3.65, 2.45, "导入数据库软件运行效果")
    add_image_contain(slide, SDK_IMG[1], 4.82, 1.72, 3.65, 2.45, "裁剪软件运行结果")
    add_image_contain(slide, SDK_IMG[2], 8.9, 1.72, 3.65, 2.45, "缩略图运行界面")
    add_image_contain(slide, SDK_IMG[3], 2.0, 4.72, 9.3, 1.55, "高清图像连续浏览截图")

    slide = new_slide("F7FAFC")
    add_brand(slide, page); page += 1
    add_title(slide, "项目一：量化指标", "代码规模来自当前仓库统计；效率为复盘估算，建议赛前按真实记录校准。")
    add_metric(slide, "13,024", "源码/工程行数", "不含 .git、.vs、OpenCV 第三方头文件和参赛资料", 0.85, 1.85, 2.35, C["blue"])
    add_metric(slide, "212", "相关文件数", "含 SDK、DemoApp、ImageSlicer 工程与源码文件", 3.55, 1.85, 2.25, C["green"])
    add_metric(slide, "85h", "节省工时估算", "传统 120h / AI 协作 35h", 6.15, 1.85, 2.25, C["yellow"])
    add_metric(slide, "71%", "效率提升估算", "以节省工时 / 传统工时计算", 8.78, 1.85, 2.25, C["red"])
    add_metric(slide, "4", "运行截图证据", "切片、导入、缩略图、高清浏览", 11.1, 1.85, 1.45, C["cyan"])
    add_panel(slide, 0.85, 3.42, 5.85, 2.65)
    add_label(slide, "模块代码分布", 1.05, 3.68, C["blue"])
    add_bars(slide, [("SDK 主体", 72, 24), ("DemoApp", 34, 7), ("ImageSlicer", 22, 4)], 1.05, 4.08, 5.3, 72, RGBColor(167, 181, 197), C["blue"])
    add_panel(slide, 7.05, 3.42, 5.5, 2.65)
    add_label(slide, "可量化能力", 7.25, 3.68, C["green"])
    add_multiline(slide, [
        "支持 GB 级超高分辨率隧道影像的切片化浏览。",
        "LOD + 动态切片加载降低内存占用和等待时间。",
        "README 明确要求 64-bit 构建，避免大图处理 2GB 内存限制。",
        "病害定位采用二分查找思路，面向图幅归属做快速定位。",
    ], 7.25, 4.1, 4.95, 1.35, size=9.2)

    slide = new_slide()
    add_brand(slide, page); page += 1
    add_title(slide, "项目一：结论及体会", "AI 的价值不只在生成代码，而在把性能策略、工程结构、调试经验和说明文档串成闭环。")
    for i, (heading, body, color, fill, x) in enumerate([
        ("结论", "已形成可演示的大图切片、缩略图浏览、高清瓦片加载和病害标注能力，可作为后续展示平台或检测业务的底层图像能力。", C["blue"], C["pale_blue"], 0.85),
        ("体会", "对 AI 说明清楚数据规模、构建环境和交互目标后，AI 更适合承担架构拆分、重复代码、问题定位和文档沉淀。", C["green"], C["pale_green"], 4.88),
        ("建议", "保留真实提示词、构建错误、运行截图和工时记录。AI 生成内容必须经过编译、真实数据和人工业务规则三重校验。", RGBColor(164, 106, 0), RGBColor(255, 247, 230), 8.9),
    ]):
        add_panel(slide, x, 1.78, 3.6, 3.55, fill)
        add_text(slide, heading, x + 0.27, 2.14, 3.0, 0.28, size=16, color=color, bold=True)
        add_text(slide, body, x + 0.27, 2.72, 3.0, 1.35, size=12, color=C["ink"])

    slide = new_slide()
    add_brand(slide, page); page += 1
    add_title(slide, "项目二：展示平台 - AI 工具与沟通过程", "从编译失败、数据映射不一致到可演示平台，AI 参与诊断、实现、验证和材料沉淀。")
    add_image_contain(slide, PLATFORM_IMG[0], 0.75, 1.75, 4.0, 2.85, "项目代码列表 / 工程规模证据")
    add_multiline(slide, [
        "Codex：阅读代码、定位编译错误、修改实体/DTO/Controller/Service、补齐文档。",
        "ChatGPT / Gemini：用于界面构思、架构讨论、表达优化和风险治理补充。",
        "沟通输入包括真实报错、业务规则、样例目录、概要设计和人工验收反馈。",
    ], 5.15, 1.82, 6.8, 1.18, size=11)
    add_table(slide, [
        ["协作阶段", "问题", "AI输出"],
        ["诊断", "新旧实体、数据库表和前端展示不一致", "影响面分析和修改清单"],
        ["实现", "台账同步、SQLite解析、统一库、接口和页面联动", "跨模块代码补齐"],
        ["验证", "构建错误、样例导入、Swagger、页面截图", "循环修复与材料沉淀"],
    ], 5.15, 3.38, [1.28, 2.65, 2.72], 0.46)
    add_placeholder(slide, 0.75, 5.12, 11.3, 1.15, "待补充：Codex / ChatGPT / Gemini 对话过程截图，可放 2-3 张关键沟通截图")

    slide = new_slide("F8FBFD")
    add_brand(slide, page); page += 1
    add_title(slide, "项目二：软件介绍与技术亮点", "展示平台面向成果汇报和客户演示，将台账、SQLite、图片和点云目录转成统一查询与可视化能力。")
    add_image_contain(slide, PLATFORM_IMG[3], 0.75, 1.72, 5.6, 2.85, "工程信息界面")
    add_image_contain(slide, PLATFORM_IMG[4], 6.75, 1.72, 5.6, 2.85, "隧道区间浏览界面")
    add_multiline(slide, [
        "链路闭环：WinForms 上传工具、统一库、WebAPI、Web 前端和 Swagger 接口。",
        "数据契约：以工程实例、站点/区间、病害、环片、灰度图和高清图组织成果。",
        "演示友好：工程地图、区间浏览、病害列表、高清图弹窗和接口文档可直接展示。",
        "风险治理：AI 生成内容以概要设计、真实表结构、构建和样例导入验证为准。",
    ], 0.9, 5.03, 11.25, 1.02, size=10.2)

    slide = new_slide()
    add_brand(slide, page); page += 1
    add_title(slide, "项目二：运行效果截图", "多端能力覆盖上传工具、Web 登录、工程信息、区间浏览和 Swagger 接口。")
    add_image_contain(slide, PLATFORM_IMG[1], 0.72, 1.68, 3.0, 1.72, "上传工具")
    add_image_contain(slide, PLATFORM_IMG[2], 4.0, 1.68, 3.0, 1.72, "登录界面")
    add_image_contain(slide, PLATFORM_IMG[3], 7.28, 1.68, 3.0, 1.72, "工程信息")
    add_image_contain(slide, PLATFORM_IMG[4], 1.1, 4.15, 5.15, 1.95, "隧道区间浏览")
    add_image_contain(slide, PLATFORM_IMG[5], 7.1, 4.15, 4.4, 1.95, "Swagger 接口")

    slide = new_slide("F7FAFC")
    add_brand(slide, page); page += 1
    add_title(slide, "项目二：量化指标", "指标来自既有《02_效果数据_量化指标.xlsx》，效率口径为复盘估算。")
    for value, label, note, x, w, color in [
        ("76%", "效率提升", "传统估算 162h / AI 协作 39h", 0.85, 2.05, C["blue"]),
        ("123h", "节省工时", "用于答辩表达，可赛前校准", 3.12, 2.05, C["green"]),
        ("23", "API接口", "Controller 中公开 HTTP 接口", 5.4, 1.65, C["cyan"]),
        ("1,461", "病害记录", "进入统一查询链路", 7.36, 1.85, C["red"]),
        ("1.8GB", "样例数据", "763 个样例文件", 9.48, 1.85, C["yellow"]),
        ("12,215", "源码行数", "不含日志和运行输出", 11.25, 1.45, C["blue"]),
    ]:
        add_metric(slide, value, label, note, x, 1.78, w, color)
    add_panel(slide, 0.85, 3.35, 6.0, 2.72)
    add_label(slide, "效率测算", 1.05, 3.62, C["blue"])
    add_bars(slide, [("需求拆解", 8, 1.5), ("数据库映射", 24, 5), ("上传工具", 20, 4.5), ("SQLite入库", 40, 11), ("前端展示", 36, 9)], 1.05, 4.02, 5.35, 40, RGBColor(167, 181, 197), C["blue"])
    add_panel(slide, 7.25, 3.35, 5.25, 2.72)
    add_label(slide, "项目数据", 7.48, 3.62, C["green"])
    add_table(slide, [["指标", "数值"], ["项目实例", "2"], ["站点/区间实体", "7"], ["已上传实体", "3"], ["环片记录", "1,258"], ["灰度图索引", "34"]], 7.48, 4.02, [2.6, 2.08], 0.32)

    slide = new_slide()
    add_brand(slide, page); page += 1
    add_title(slide, "项目二：结论及体会", "AI 加速了从业务资料到可运行系统的闭环，但核心规则仍需人工验收。")
    add_multiline(slide, [
        "结论：展示平台已形成上传、入库、接口、页面和材料沉淀的演示闭环。",
        "体会：AI 最适合承担跨模块定位、样板代码生成、文档整理和反复修复。",
        "边界：涉及 SQLite 字段含义、Word 表语义、工程数据安全时，必须以真实数据和人工规则为准。",
        "推广：适合复制到“多源文件 + 结构化入库 + 可视化展示 + 汇报沉淀”的业务场景。",
    ], 0.9, 1.85, 6.0, 2.0, size=13)
    add_table(slide, [
        ["AI风险", "控制措施", "竞赛表达"],
        ["幻觉/误映射", "以概要设计和真实表结构验证", "AI加速，人工把关"],
        ["数据安全", "本地资料、本地运行", "适合内网部署"],
        ["质量一致性", "构建、接口、页面检查", "形成可复跑清单"],
        ["知识沉淀", "报告、日志、API说明", "个人试错转为团队资产"],
    ], 7.15, 1.9, [1.18, 2.05, 2.02], 0.47)

    slide = new_slide("F3F7FA")
    add_brand(slide, page); page += 1
    add_title(slide, "综合结论：AI 已经进入研发交付主流程", "两项项目证明，AI 能把需求拆解、代码生成、调试验证和竞赛材料沉淀联成一套方法。")
    for heading, body, color, x in [
        ("可复制方法", "用真实业务资料驱动 AI，不只问“怎么写代码”，而是持续提供报错、截图、数据结构和验收结果。", C["blue"], 0.85),
        ("可量化价值", "SDK：约 71% 效率提升估算；展示平台：约 76% 效率提升估算。两者均需保留真实工时记录持续校准。", C["green"], 4.95),
        ("后续建议", "补齐 AI 对话截图、真实工时台账、演示录屏和脱敏检查，把本次竞赛材料升级为部门 AI 研发实践模板。", RGBColor(164, 106, 0), 9.05),
    ]:
        add_panel(slide, x, 1.85, 3.35, 3.55)
        add_text(slide, heading, x + 0.3, 2.18, 2.6, 0.3, size=16, bold=True, color=color)
        add_text(slide, body, x + 0.3, 2.82, 2.68, 1.25, size=12, color=C["ink"])
    add_text(slide, "建议答辩主线：痛点 → AI 协作方法 → 两个项目成果 → 数据指标 → 风险治理 → 推广价值", 0.92, 6.0, 11.5, 0.36, size=13, bold=True, color=C["navy"], align=PP_ALIGN.CENTER)


def write_previews():
    try:
        font_big = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 54)
        font_mid = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 28)
        font_num = ImageFont.truetype("C:/Windows/Fonts/arial.ttf", 36)
    except Exception:
        font_big = font_mid = font_num = None
    previews = [
        ("01", "AI 赋能研发交付", "大图连续浏览 SDK 与地铁隧道展示平台"),
        ("02", "两项 AI 实践形成一条可复用链路", "切片 / 浏览标注 / 统一库接口 / 可视化汇报"),
        ("03", "项目一：AI 工具与沟通过程", "保留 AI 对话截图占位"),
        ("04", "项目一：软件介绍与技术亮点", "LOD、动态瓦片、异步加载、专业标注"),
        ("05", "项目一：运行效果截图", "切片、缩略图、高清浏览"),
        ("06", "项目一：量化指标", "13,024 行，约 71% 效率提升估算"),
        ("07", "项目一：结论及体会", "AI 串联性能策略、工程结构和文档"),
        ("08", "项目二：AI 工具与沟通过程", "从编译失败到可演示系统"),
        ("09", "项目二：软件介绍与技术亮点", "统一库、接口、Web 展示"),
        ("10", "项目二：运行效果截图", "上传工具、页面、Swagger"),
        ("11", "项目二：量化指标", "76% 效率提升，节省 123h"),
        ("12", "项目二：结论及体会", "AI 加速，人工验收把关"),
        ("13", "综合结论", "AI 已进入研发交付主流程"),
    ]
    for no, title, sub in previews:
        im = Image.new("RGB", (1600, 900), (247, 250, 252))
        draw = ImageDraw.Draw(im)
        draw.rounded_rectangle((74, 76, 1526, 824), radius=22, fill=(255, 255, 255), outline=(215, 225, 234), width=2)
        draw.text((108, 110), no, fill=(29, 93, 155), font=font_num)
        draw.text((108, 245), title, fill=(24, 33, 47), font=font_big)
        draw.text((110, 325), sub, fill=(95, 107, 122), font=font_mid)
        draw.rectangle((110, 620, 1090, 634), fill=(29, 93, 155))
        draw.text((110, 690), "AI竞赛综合汇报 PPT 预览图，用于快速核对页序和主题。", fill=(95, 107, 122), font=font_mid)
        im.save(PREVIEW_DIR / f"slide-{no}.png")


if __name__ == "__main__":
    build()
    prs.save(OUT)
    write_previews()
    print(OUT)
    print(PREVIEW_DIR)
